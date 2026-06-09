from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path("docs/business/frontend_client_ui_requirements_20260416.pptx")

BG = RGBColor(246, 244, 239)
NAVY = RGBColor(28, 44, 67)
TEAL = RGBColor(42, 122, 120)
ORANGE = RGBColor(211, 120, 56)
TEXT = RGBColor(44, 52, 64)
MUTED = RGBColor(107, 114, 128)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 214, 208)


def add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.0), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = NAVY
    if subtitle:
        sp = tf.add_paragraph()
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Aptos"
        sr.font.size = Pt(10)
        sr.font.color.rgb = MUTED
        sp.space_before = Pt(4)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.35), Inches(2.1), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_bullets(slide, left, top, width, height, bullets: list[str], font_size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
        p.bullet = True
    return box


def add_card(slide, left, top, width, height, title: str, bullets: list[str], accent=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.12), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.12), width - Inches(0.3), Inches(0.35))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = NAVY

    add_bullets(slide, left + Inches(0.18), top + Inches(0.48), width - Inches(0.3), height - Inches(0.55), bullets, 13)


def add_status_flow(slide, labels: list[str]):
    start_left = 0.8
    top = 2.2
    width = 2.1
    gap = 0.25
    for idx, label in enumerate(labels):
        left = start_left + idx * (width + gap)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.85)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = TEAL if idx in {1, 2, 3} else LINE
        shape.line.width = Pt(1.25)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = TEXT
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + width - 0.02),
                Inches(top + 0.2),
                Inches(0.32),
                Inches(0.42),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE
            arrow.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = [
        {
            "title": "面對用戶端前端介面需求草案",
            "subtitle": "根據目前需求整理的第一版討論稿",
            "cards": [
                ("這份簡報回答什麼", [
                    "前端至少要處理哪些角色、權限與流程",
                    "目前想到的需求還缺哪些關鍵面向",
                    "先做 MVP 時，最少需要哪些畫面與資料",
                ], TEAL),
                ("目前已知核心需求", [
                    "項目需要新 / 刪 / 修權限",
                    "不同角色有不同操作行為",
                    "缺料或找不到物料時要能提需求",
                    "新增或刪除物料需求需要審核與時間紀錄",
                ], ORANGE),
            ],
        },
        {
            "title": "系統核心物件",
            "subtitle": "先分清楚人、項目、物料、需求、審核五個層次",
            "cards": [
                ("User / Role", ["使用者", "角色", "部門或資料範圍", "帳號狀態"], TEAL),
                ("Project", ["項目主檔", "項目狀態", "建立人 / 維護人", "最後更新時間"], TEAL),
                ("Material", ["物料主檔", "物料代碼", "規格 / 單位", "啟用狀態"], TEAL),
                ("Request", ["新增物料需求", "刪除物料需求", "修改物料需求", "送審狀態"], ORANGE),
                ("Approval", ["審核人", "審核結果", "審核意見", "審核時間"], ORANGE),
            ],
        },
        {
            "title": "角色與權限設計",
            "subtitle": "角色不等於權限，還要考慮資料範圍",
            "cards": [
                ("角色建議", [
                    "一般使用者：查看項目、補缺少品項、提需求",
                    "項目維護者：可直接編輯項目內物料",
                    "審核者：核准或駁回新增 / 刪除需求",
                    "管理者：管理角色、權限、主檔資料",
                ], TEAL),
                ("權限建議", [
                    "View：查看",
                    "Create：新增",
                    "Edit：修改",
                    "Delete：刪除或作廢",
                    "Submit / Approve：送審與審核",
                ], ORANGE),
            ],
        },
        {
            "title": "你目前比較容易漏掉的點",
            "subtitle": "如果現在不補，前端做一半通常會回頭重改",
            "cards": [
                ("流程面", ["需求狀態", "駁回後能否重送", "已核准後能否撤回", "刪除是否為真刪除"], TEAL),
                ("資料面", ["建立時間", "送審時間", "審核時間", "最後更新時間", "操作人"], ORANGE),
                ("控制面", ["部門 / 專案資料範圍", "並發修改衝突", "通知機制", "附件與備註"], TEAL),
            ],
        },
        {
            "title": "建議的狀態流轉",
            "subtitle": "新增 / 刪除 / 修改物料需求都建議走同一套審核狀態",
            "flow": ["草稿", "已送審", "審核中", "已通過", "已駁回", "已取消"],
            "bullets": [
                "所有需求都應保留申請人、審核人、審核時間與意見",
                "駁回不應直接消失，應能回看歷程與重新送審",
                "已通過後若涉及正式資料變更，應再寫入操作紀錄",
            ],
        },
        {
            "title": "前端畫面最少需要哪些",
            "subtitle": "先做 MVP 時，這些畫面通常跑不掉",
            "cards": [
                ("核心操作頁", [
                    "登入頁",
                    "項目列表",
                    "項目詳情",
                    "項目內物料清單",
                    "新增 / 編輯物料視窗",
                ], TEAL),
                ("流程支援頁", [
                    "找不到物料時的需求申請視窗",
                    "審核待辦列表",
                    "審核詳情頁",
                    "歷程 / 操作紀錄區塊",
                ], ORANGE),
            ],
        },
        {
            "title": "每筆需求至少要存的欄位",
            "subtitle": "這頁可以直接變成後端與前端欄位清單的起點",
            "cards": [
                ("需求主欄位", [
                    "request_id",
                    "project_id",
                    "request_type",
                    "material_id 或臨時名稱",
                    "reason / remark",
                    "status",
                ], TEAL),
                ("時間與人員欄位", [
                    "requested_by",
                    "requested_at",
                    "submitted_at",
                    "approved_by",
                    "approved_at",
                    "updated_at",
                ], ORANGE),
            ],
        },
        {
            "title": "我建議先做的 MVP 範圍",
            "subtitle": "先讓流程能跑，再補通知、報表、細緻規則",
            "cards": [
                ("MVP 要先有", [
                    "登入",
                    "角色與基本權限",
                    "項目內物料增修",
                    "找不到物料時可提需求",
                    "審核通過 / 駁回",
                    "時間與操作紀錄",
                ], TEAL),
                ("第二階段再做", [
                    "通知中心",
                    "附件上傳",
                    "細緻部門權限",
                    "批次匯入匯出",
                    "報表與統計看板",
                ], ORANGE),
            ],
        },
        {
            "title": "下一步建議",
            "subtitle": "這份簡報之後最適合接的工作順序",
            "cards": [
                ("先定義", [
                    "角色 x 權限表",
                    "需求狀態流轉表",
                    "每個畫面的欄位與操作",
                ], TEAL),
                ("再產出", [
                    "前端 sitemap",
                    "low-fidelity wireframe",
                    "API 與資料模型草稿",
                ], ORANGE),
            ],
        },
    ]

    for spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)
        add_title(slide, spec["title"], spec.get("subtitle"))
        if "flow" in spec:
            add_status_flow(slide, spec["flow"])
            add_bullets(slide, Inches(0.9), Inches(3.5), Inches(11.5), Inches(2.3), spec["bullets"], 17)
            continue

        cards = spec["cards"]
        count = len(cards)
        if count == 2:
            positions = [
                (Inches(0.85), Inches(1.75), Inches(5.75), Inches(4.8)),
                (Inches(6.73), Inches(1.75), Inches(5.75), Inches(4.8)),
            ]
        elif count == 3:
            positions = [
                (Inches(0.65), Inches(1.75), Inches(4.0), Inches(4.8)),
                (Inches(4.66), Inches(1.75), Inches(4.0), Inches(4.8)),
                (Inches(8.67), Inches(1.75), Inches(4.0), Inches(4.8)),
            ]
        else:
            positions = [
                (Inches(0.45), Inches(1.75), Inches(2.45), Inches(4.8)),
                (Inches(2.95), Inches(1.75), Inches(2.45), Inches(4.8)),
                (Inches(5.45), Inches(1.75), Inches(2.45), Inches(4.8)),
                (Inches(7.95), Inches(1.75), Inches(2.45), Inches(4.8)),
                (Inches(10.45), Inches(1.75), Inches(2.45), Inches(4.8)),
            ]

        for idx, (title, bullets, accent) in enumerate(cards):
            left, top, width, height = positions[idx]
            add_card(slide, left, top, width, height, title, bullets, accent)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH.resolve())


if __name__ == "__main__":
    build_presentation()
